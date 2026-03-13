# main.py - FastAPI Backend for HeyGen LiveAvatar Streaming + Document Ingestion Pipeline
import asyncio
import json
import logging
import shutil
import time
import uuid
from pathlib import Path

from fastapi import FastAPI, File, Request, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
import httpx
import os
from dotenv import load_dotenv
from openai import AsyncOpenAI

load_dotenv()

HEYGEN_API_KEY = os.getenv("HEYGEN_API_KEY")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
BASE_URL = "https://api.heygen.com"
RAG_API_URL = "https://rag-super-agent.onrender.com/chat/"

# Storage directories
UPLOAD_DIR = Path("data/uploads")
OUTPUT_DIR = Path("data/outputs")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# In-memory job tracker (no Celery/Redis needed)
jobs: dict = {}

SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".doc", ".docx", ".pptx", ".csv", ".md", ".rtf"}

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(title="HeyGen LiveAvatar + Ingestion Pipeline", version="2.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ─── OpenAI client ───────────────────────────────────────────────────────────
openai_client = AsyncOpenAI(api_key=OPENAI_API_KEY) if OPENAI_API_KEY else None

# ─── Text Extraction Helpers ─────────────────────────────────────────────────

def extract_text(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(file_path)
    elif ext in (".txt", ".md"):
        return _extract_plain(file_path)
    elif ext == ".docx":
        return _extract_docx(file_path)
    elif ext == ".doc":
        return _extract_doc(file_path)
    elif ext == ".pptx":
        return _extract_pptx(file_path)
    elif ext == ".csv":
        return _extract_csv(file_path)
    elif ext == ".rtf":
        return _extract_rtf(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def _extract_pdf(path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(path)
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_plain(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _extract_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)


def _extract_doc(path: str) -> str:
    import subprocess
    for cmd in ["antiword", "catdoc"]:
        try:
            result = subprocess.run([cmd, path], capture_output=True, text=True, timeout=30)
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout
        except FileNotFoundError:
            continue
    raise RuntimeError("Cannot extract .doc: install antiword or catdoc, or convert to .docx")


def _extract_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    texts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
        if slide_texts:
            texts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
    return "\n\n".join(texts)


def _extract_csv(path: str) -> str:
    import csv as csv_mod
    rows = []
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        for row in csv_mod.reader(f):
            rows.append(" | ".join(row))
    return "\n".join(rows)


def _extract_rtf(path: str) -> str:
    import re
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()
    text = re.sub(r'\\[a-z]+\d*\s?', '', content)
    text = re.sub(r'[{}]', '', text)
    return text.strip()


# ─── LLM Script Generation ──────────────────────────────────────────────────

async def generate_video_script(text: str, custom_prompt: str = "") -> str:
    if not openai_client:
        raise RuntimeError("OPENAI_API_KEY not set")
        
    system_prompt = custom_prompt.strip() if custom_prompt else (
        "You are a professional video script writer. Given document content, "
        "write a concise, engaging narration script (60-90 seconds when spoken) "
        "suitable for an AI avatar video. The script should summarize the key "
        "points clearly and engagingly. Output ONLY the script text, no stage "
        "directions or formatting."
    )

    # Use first 4000 chars to stay within token limits
    snippet = text[:4000]
    response = await openai_client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {
                "role": "system",
                "content": system_prompt,
            },
            {
                "role": "user",
                "content": f"Write a narration script based on this document content:\n\n{snippet}",
            },
        ],
        max_tokens=1000,
        temperature=0.7,
    )
    return response.choices[0].message.content.strip()


# ─── HeyGen Video Generation ────────────────────────────────────────────────

async def create_heygen_video(script: str) -> str | None:
    """Submit video generation to HeyGen, return video_id."""
    if not HEYGEN_API_KEY:
        raise RuntimeError("HEYGEN_API_KEY not set")

    headers = {"X-Api-Key": HEYGEN_API_KEY, "Content-Type": "application/json"}

    # Fetch first available avatar
    avatar_id = "Anna_public_3_20240108"
    voice_id = "1bd001e7e50f421d891986aad5158bc8"
    try:
        async with httpx.AsyncClient(timeout=30) as client:
            res = await client.get(f"{BASE_URL}/v2/avatars", headers=headers)
            if res.status_code == 200:
                avatars = res.json().get("data", {}).get("avatars", [])
                if avatars:
                    avatar_id = avatars[0].get("avatar_id", avatar_id)
    except Exception as e:
        logger.warning(f"Could not fetch avatars, using default: {e}")

    payload = {
        "video_inputs": [
            {
                "character": {"type": "avatar", "avatar_id": avatar_id, "avatar_style": "normal"},
                "voice": {"type": "text", "input_text": script, "voice_id": voice_id},
            }
        ],
        "dimension": {"width": 1280, "height": 720},
    }

    async with httpx.AsyncClient(timeout=60) as client:
        res = await client.post(f"{BASE_URL}/v2/video/generate", json=payload, headers=headers)
        res.raise_for_status()
        data = res.json()
        return data.get("data", {}).get("video_id")


async def poll_heygen_video(video_id: str, job_id: str):
    """Poll HeyGen for video completion and download the MP4."""
    headers = {"X-Api-Key": HEYGEN_API_KEY}
    status_url = f"{BASE_URL}/v1/video_status.get?video_id={video_id}"

    for attempt in range(60):
        try:
            async with httpx.AsyncClient(timeout=30) as client:
                res = await client.get(status_url, headers=headers)
                data = res.json()

            status = data.get("data", {}).get("status")
            jobs[job_id]["heygen_status"] = status

            if status == "completed":
                video_url = data.get("data", {}).get("video_url")
                if video_url:
                    # Download MP4
                    mp4_path = OUTPUT_DIR / f"{job_id}_video.mp4"
                    async with httpx.AsyncClient(timeout=120) as client:
                        dl = await client.get(video_url)
                        with open(mp4_path, "wb") as f:
                            f.write(dl.content)
                    jobs[job_id]["stage"] = "completed"
                    jobs[job_id]["video_file"] = str(mp4_path)
                    jobs[job_id]["video_url"] = f"/download/{job_id}"
                    logger.info(f"Job {job_id}: Video downloaded to {mp4_path}")
                return
            elif status == "failed":
                jobs[job_id]["stage"] = "failed"
                error_info = data.get("data", {}).get("error", {})
                error_msg = error_info.get("message", error_info.get("detail", "HeyGen video generation failed"))
                jobs[job_id]["error"] = f"HeyGen Error: {error_msg}"
                return
            else:
                jobs[job_id]["message"] = f"HeyGen rendering... ({status})"
                await asyncio.sleep(10)
        except Exception as e:
            logger.error(f"Poll error: {e}")
            await asyncio.sleep(10)

    jobs[job_id]["stage"] = "failed"
    jobs[job_id]["error"] = "HeyGen video timed out after 10 minutes"


# ─── Background Pipeline ────────────────────────────────────────────────────

async def run_ingestion_pipeline(job_id: str, file_path: str, filename: str):
    """Full pipeline: extract text -> LLM script -> HeyGen video."""
    try:
        # Stage 1: Extract text
        jobs[job_id]["stage"] = "extracting"
        jobs[job_id]["message"] = "Extracting text from document..."
        jobs[job_id]["step"] = 1
        text = extract_text(file_path)
        jobs[job_id]["extracted_length"] = len(text)
        logger.info(f"Job {job_id}: Extracted {len(text)} chars from {filename}")

        # Save extracted text
        text_path = OUTPUT_DIR / f"{job_id}_text.txt"
        with open(text_path, "w", encoding="utf-8") as f:
            f.write(text)

        # Stage 2: Wait for user to review text and provide prompt
        jobs[job_id]["extracted_text"] = text
        jobs[job_id]["stage"] = "text_extracted"
        jobs[job_id]["message"] = "Text extracted. Ready for AI Prompt."
        jobs[job_id]["step"] = 2

    except Exception as e:
        logger.error(f"Job {job_id} failed: {e}", exc_info=True)
        jobs[job_id]["stage"] = "failed"
        jobs[job_id]["error"] = str(e)


async def run_script_generation(job_id: str, prompt: str):
    """Phase 2 pipeline: submit extracted text and prompt to LLM."""
    try:
        jobs[job_id]["stage"] = "generating_script"
        jobs[job_id]["message"] = "Generating video script with AI..."
        jobs[job_id]["step"] = 3
        
        text = jobs[job_id].get("extracted_text", "")
        script = await generate_video_script(text, prompt)
        jobs[job_id]["script"] = script
        logger.info(f"Job {job_id}: Script generated ({len(script)} chars)")

        # Save script
        script_path = OUTPUT_DIR / f"{job_id}_script.txt"
        with open(script_path, "w", encoding="utf-8") as f:
            f.write(script)

        # Stage 3: Wait for user to review/edit
        jobs[job_id]["stage"] = "script_ready"
        jobs[job_id]["message"] = "Script is ready for review and editing."
        jobs[job_id]["step"] = 4

    except Exception as e:
        logger.error(f"Job {job_id} failed: {e}", exc_info=True)
        jobs[job_id]["stage"] = "failed"
        jobs[job_id]["error"] = str(e)


async def run_video_generation(job_id: str, script: str):
    """Phase 3 pipeline: submit script to HeyGen -> render video."""
    try:
        # Update script in memory
        jobs[job_id]["script"] = script

        # Stage 4: Submit to HeyGen
        jobs[job_id]["stage"] = "submitting_video"
        jobs[job_id]["message"] = "Submitting to HeyGen for video generation..."
        jobs[job_id]["step"] = 5
        video_id = await create_heygen_video(script)

        if not video_id:
            jobs[job_id]["stage"] = "failed"
            jobs[job_id]["error"] = "HeyGen did not return a video_id"
            return

        jobs[job_id]["heygen_video_id"] = video_id

        # Stage 5: Poll for completion
        jobs[job_id]["stage"] = "rendering"
        jobs[job_id]["message"] = "HeyGen is rendering the video..."
        jobs[job_id]["step"] = 6
        await poll_heygen_video(video_id, job_id)

    except Exception as e:
        logger.error(f"Job {job_id} failed: {e}", exc_info=True)
        jobs[job_id]["stage"] = "failed"
        jobs[job_id]["error"] = str(e)


# ═══════════════════════════════════════════════════════════════════════════════
#  EXISTING STREAMING ENDPOINTS (unchanged)
# ═══════════════════════════════════════════════════════════════════════════════

async def heygen_get(endpoint: str):
    async with httpx.AsyncClient(timeout=60.0) as client:
        res = await client.get(f"{BASE_URL}/{endpoint}", headers={"X-Api-Key": HEYGEN_API_KEY})
        res.raise_for_status()
        return res.json()


async def heygen_post(endpoint: str, data: dict = None):
    async with httpx.AsyncClient(timeout=60.0) as client:
        res = await client.post(
            f"{BASE_URL}/{endpoint}",
            json=data,
            headers={"X-Api-Key": HEYGEN_API_KEY, "Content-Type": "application/json"},
        )
        res.raise_for_status()
        return res.json()


@app.get("/")
async def root():
    return {
        "status": "ok",
        "service": "HeyGen LiveAvatar + Ingestion Pipeline",
        "endpoints": {
            "streaming_token": "/streaming_token (POST)",
            "avatars": "/avatars",
            "interactive_avatars": "/interactive_avatars",
            "voices": "/voices",
            "chat": "/chat (POST)",
            "ingest": "/ingest (POST - file upload)",
            "job_status": "/job/{job_id} (GET)",
            "download": "/download/{job_id} (GET)",
        },
    }


@app.post("/streaming_token")
async def get_streaming_token():
    try:
        async with httpx.AsyncClient(timeout=60.0) as client:
            res = await client.post(
                f"{BASE_URL}/v1/streaming.create_token",
                headers={"X-Api-Key": HEYGEN_API_KEY, "Content-Type": "application/json"},
            )
            if res.status_code != 200:
                return JSONResponse(status_code=res.status_code, content={"error": res.text})
            return res.json()
    except httpx.TimeoutException as e:
        return JSONResponse(status_code=504, content={"error": f"Request timeout: {e}"})
    except Exception as e:
        import traceback
        return JSONResponse(status_code=500, content={"error": str(e), "traceback": traceback.format_exc()})


@app.get("/avatars")
async def get_avatars():
    try:
        return await heygen_get("v2/avatars")
    except httpx.TimeoutException as e:
        return JSONResponse(status_code=504, content={"error": f"Request timeout: {e}"})
    except httpx.HTTPStatusError as e:
        return JSONResponse(status_code=e.response.status_code, content={"error": e.response.text})
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})


@app.get("/interactive_avatars")
async def get_interactive_avatars():
    try:
        return await heygen_get("v1/streaming/avatar.list")
    except httpx.TimeoutException as e:
        return JSONResponse(status_code=504, content={"error": f"Request timeout: {e}"})
    except httpx.HTTPStatusError as e:
        return JSONResponse(status_code=e.response.status_code, content={"error": e.response.text})
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})


@app.get("/voices")
async def get_voices():
    try:
        return await heygen_get("v2/voices")
    except httpx.TimeoutException as e:
        return JSONResponse(status_code=504, content={"error": f"Request timeout: {e}"})
    except httpx.HTTPStatusError as e:
        return JSONResponse(status_code=e.response.status_code, content={"error": e.response.text})
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})


@app.post("/chat")
async def chat(request: Request):
    try:
        data = await request.json()
        message = data.get("message")
        if not message:
            return JSONResponse(status_code=400, content={"error": "Message is required"})

        payload = {"message": message}
        if data.get("conversation_id"):
            payload["conversation_id"] = data["conversation_id"]

        async with httpx.AsyncClient(timeout=120.0) as client:
            res = await client.post(RAG_API_URL, json=payload, headers={"Content-Type": "application/json"})
            res.raise_for_status()
            return res.json()
    except httpx.TimeoutException as e:
        return JSONResponse(status_code=504, content={"error": f"Request timeout: {e}"})
    except httpx.HTTPStatusError as e:
        return JSONResponse(status_code=e.response.status_code, content={"error": e.response.text})
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})


# ═══════════════════════════════════════════════════════════════════════════════
#  NEW INGESTION PIPELINE ENDPOINTS
# ═══════════════════════════════════════════════════════════════════════════════

@app.post("/ingest")
async def ingest_document(file: UploadFile = File(...)):
    """
    Upload a document (PDF, DOCX, TXT, PPTX, CSV, MD, RTF),
    extract text, generate script via LLM, create HeyGen video.
    Returns a job_id to poll for status.
    """
    ext = Path(file.filename).suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        return JSONResponse(
            status_code=400,
            content={"error": f"Unsupported file type: {ext}. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"},
        )

    job_id = str(uuid.uuid4())
    file_path = UPLOAD_DIR / f"{job_id}_{file.filename}"

    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    jobs[job_id] = {
        "job_id": job_id,
        "filename": file.filename,
        "stage": "uploaded",
        "message": "File uploaded, starting pipeline...",
        "step": 0,
        "total_steps": 4,
        "script": None,
        "video_file": None,
        "video_url": None,
        "error": None,
    }

    # Run phase 1 pipeline in background
    asyncio.create_task(run_ingestion_pipeline(job_id, str(file_path), file.filename))

    return {
        "message": "Ingestion started",
        "job_id": job_id,
        "filename": file.filename,
        "status_endpoint": f"/job/{job_id}",
    }


@app.post("/job/{job_id}/generate_script")
async def generate_script(job_id: str, request: Request):
    """
    Phase 2: User provides custom prompt for the AI script generation.
    """
    job = jobs.get(job_id)
    if not job:
        return JSONResponse(status_code=404, content={"error": "Job not found"})
        
    data = await request.json()
    prompt = data.get("prompt", "")
    
    if job["stage"] not in ["text_extracted", "failed"]:
        return JSONResponse(
            status_code=400, 
            content={"error": f"Job is in stage '{job['stage']}', cannot generate script now"}
        )

    # Move to phase 2
    job["stage"] = "generating_script"
    job["message"] = "Generating script using custom prompt..."
    job["step"] = 3
    job["error"] = None

    asyncio.create_task(run_script_generation(job_id, prompt))

    return {"message": "Script generation started"}


@app.post("/job/{job_id}/generate_video")
async def generate_video_from_script(job_id: str, request: Request):
    """
    Phase 3: User reviewed and edited the script. 
    Submit the edited script to HeyGen for video generation.
    """
    job = jobs.get(job_id)
    if not job:
        return JSONResponse(status_code=404, content={"error": "Job not found"})
        
    data = await request.json()
    script = data.get("script")
    
    if not script:
        return JSONResponse(status_code=400, content={"error": "Script text is required"})

    if job["stage"] not in ["script_ready", "failed"]:
        return JSONResponse(
            status_code=400, 
            content={"error": f"Job is in stage '{job['stage']}', cannot generate video now"}
        )

    # Move to phase 3
    job["stage"] = "submitting_video"
    job["message"] = "Preparing to submit edited script to HeyGen..."
    job["step"] = 5
    job["error"] = None

    asyncio.create_task(run_video_generation(job_id, script))

    return {"message": "Video generation started"}


@app.get("/job/{job_id}")
async def get_job_status(job_id: str):
    """Poll job status."""
    job = jobs.get(job_id)
    if not job:
        return JSONResponse(status_code=404, content={"error": "Job not found"})
    return job


@app.get("/download/{job_id}")
async def download_video(job_id: str):
    """Download the generated MP4 video."""
    job = jobs.get(job_id)
    if not job:
        return JSONResponse(status_code=404, content={"error": "Job not found"})

    video_file = job.get("video_file")
    if not video_file or not Path(video_file).exists():
        return JSONResponse(status_code=404, content={"error": "Video not ready yet"})

    return FileResponse(
        video_file,
        media_type="video/mp4",
        filename=f"{job.get('filename', 'video')}.mp4",
    )


@app.get("/jobs")
async def list_jobs():
    """List all jobs with their current status."""
    return {"jobs": list(jobs.values())}


# ─── SPA Serving ─────────────────────────────────────────────────────────────

dist_path = Path(__file__).parent / "dist"
if dist_path.exists():
    assets_path = dist_path / "assets"
    if assets_path.exists():
        app.mount("/assets", StaticFiles(directory=assets_path), name="assets")

    @app.get("/{full_path:path}")
    async def serve_spa(full_path: str, request: Request):
        api_prefixes = (
            "streaming_token", "avatars", "interactive_avatars", "voices",
            "chat", "assets", "docs", "openapi.json", "ingest", "job", "jobs", "download",
        )
        if full_path.startswith(api_prefixes):
            return JSONResponse(status_code=404, content={"error": "Not found"})
        index_file = dist_path / "index.html"
        if index_file.exists():
            return FileResponse(index_file)
        return JSONResponse(status_code=404, content={"error": "Frontend not built"})


if __name__ == "__main__":
    import uvicorn
    port = int(os.getenv("PORT", 3002))
    uvicorn.run(app, host="0.0.0.0", port=port, reload=False)
